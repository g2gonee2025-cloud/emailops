"""
Text Extraction.

Wraps Unstructured, Tesseract OCR, pdfplumber, etc.
Implements §6.5 of the Canonical Blueprint.
"""

from __future__ import annotations

import contextlib
import csv
import io
import ipaddress
import logging
import os
import re
import threading
import time
from collections.abc import Callable, Iterable, Sequence
from pathlib import Path
from typing import Any, cast
from urllib.parse import urlparse

from cortex.config.loader import get_config
from cortex.utils import strip_control_chars

logger = logging.getLogger(__name__)

# File type support
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".log",
    ".json",
    ".yaml",
    ".yml",
    ".csv",
    ".xml",
    ".html",
    ".htm",
}
DOCX_EXTENSIONS = {".docx", ".doc"}
PDF_EXTENSIONS = {".pdf"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
PPT_EXTENSIONS = {".pptx", ".ppt"}
RTF_EXTENSIONS = {".rtf"}
EMAIL_EXTENSIONS = {".eml", ".msg"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}
MBOX_EXTENSIONS = {".mbox"}

# Global cache for expensive text extraction operations
# (cached_at, mtime, size, text)
_extraction_cache: dict[tuple[Path, int | None], tuple[float, float, int, str]] = {}
_extraction_cache_lock = threading.Lock()
_CACHE_TTL = int(os.getenv("EXTRACTION_CACHE_TTL", "3600"))  # 1 hour TTL


def _is_cache_valid(
    cache_entry: tuple[float, float, int, str],
    *,
    path: Path,
) -> bool:
    """
    Check if a cached extraction result is still valid.
    Validity requires:
      - TTL not expired, and
      - file mtime/size unchanged since cache was written.
    """
    try:
        cached_at, cached_mtime, cached_size, _txt = cache_entry
    except Exception:
        return False

    # TTL check
    if (time.time() - float(cached_at)) > _CACHE_TTL:
        return False

    try:
        st = path.stat()
        cur_mtime = float(st.st_mtime)
        cur_size = int(st.st_size)
    except Exception:
        return False

    # Use small epsilon for float comparison safety
    return (abs(cur_mtime - float(cached_mtime)) < 1e-6) and (
        cur_size == int(cached_size)
    )


def _finalize_text(text: str, max_chars: int | None) -> str:
    """Normalize, truncate, and sanitize extracted text."""
    text = text or ""
    if max_chars is not None and len(text) > max_chars:
        text = text[:max_chars]
    return strip_control_chars(text)


def _is_safe_tika_url(url: str) -> bool:
    """Validate Tika endpoint to reduce SSRF exposure."""
    try:
        parsed = urlparse(url)
    except ValueError:
        return False
    if parsed.scheme not in {"http", "https"}:
        return False
    hostname = parsed.hostname
    if not hostname:
        return False
    host_lower = hostname.lower()
    if host_lower in {"localhost", "127.0.0.1", "::1"}:
        return False
    if host_lower.endswith((".local", ".internal")):
        return False
    try:
        ip_addr = ipaddress.ip_address(hostname)
    except ValueError:
        return True
    return not (
        ip_addr.is_private
        or ip_addr.is_loopback
        or ip_addr.is_link_local
        or ip_addr.is_reserved
        or ip_addr.is_multicast
    )


def _extract_with_tika(path: Path, max_chars: int | None) -> str:
    """Extract text using Apache Tika if available."""
    try:
        from tika import parser  # type: ignore

        config = get_config()
        tika_server_url = (
            str(config.system.tika_server_endpoint)
            if config.system.tika_server_endpoint
            else os.getenv("TIKA_SERVER_URL")
        )
        # Validate URL to prevent arbitrary SSRF if variable is compromised
        if tika_server_url:
            if not _is_safe_tika_url(tika_server_url):
                logger.warning("Unsafe TIKA_SERVER_URL ignored: %s", tika_server_url)
            else:
                os.environ.setdefault("TIKA_SERVER_ENDPOINT", tika_server_url)
    except ImportError:
        logger.info("Apache Tika not installed; skipping %s", path)
        return ""
    except Exception as e:
        logger.debug("Tika import/setup failed path=%s: %s", path, e)
        return ""

    try:
        # Tika parser can throw various connection errors
        parsed_raw = parser.from_file(str(path))  # type: ignore[attr-defined]
        if isinstance(parsed_raw, dict):
            parsed_dict = cast(dict[str, Any], parsed_raw)
            content = cast(str | None, parsed_dict.get("content"))
            if content:
                return _finalize_text(content, max_chars)
    except Exception as exc:
        logger.warning("Tika parse failed for %s: %s", path, exc)
    logger.debug("Apache Tika produced no content for %s", path)
    return ""


def _extract_with_unstructured(path: Path, suffix: str, max_chars: int | None) -> str:
    """Extract text using Unstructured's partitioners when installed."""
    try:
        from unstructured.partition.auto import partition  # type: ignore
    except ImportError:
        logger.info("unstructured not installed; skipping %s", path)
        return ""
    except Exception as e:
        logger.warning("Unstructured import failed for %s: %s", path, e)
        return ""

    kwargs: dict[str, object] = {
        "filename": str(path),
        "include_page_breaks": True,
        "metadata_filename": True,
        "process_attachments": suffix in EMAIL_EXTENSIONS,
    }

    try:
        elements = cast(Iterable[Any], partition(**kwargs))
    except (ValueError, OSError, RuntimeError) as exc:
        # Catch specific known runtime errors from unstructured
        logger.warning("unstructured partition failed for %s: %s", path, exc)
        return ""
    except Exception as exc:
        # Catch-all for other unexpected unstructured errors
        logger.warning("Unexpected unstructured error for %s: %s", path, exc)
        return ""

    parts: list[str] = []
    for element in elements:
        text = cast(str | None, getattr(element, "text", None))
        if text and text.strip():
            parts.append(text.strip())

    if not parts:
        return ""

    return _finalize_text("\n\n".join(parts), max_chars)


def _extract_image_ocr(path: Path, max_chars: int | None) -> str:
    """Attempt OCR extraction for image-based attachments."""
    try:
        from PIL import Image  # type: ignore
    except ImportError:
        logger.info("Pillow not installed; skipping image OCR for %s", path)
        return ""
    except Exception as exc:
        logger.warning("Failed to import Pillow for %s: %s", path, exc)
        return ""

    try:
        import pytesseract  # type: ignore
    except ImportError:
        logger.info("pytesseract not installed; skipping OCR for %s", path)
        return ""
    except Exception as exc:
        logger.warning("Failed to import pytesseract for %s: %s", path, exc)
        return ""

    try:
        with Image.open(path) as image:
            pytesseract_module = cast(Any, pytesseract)
            text = cast(str, pytesseract_module.image_to_string(image))
    except Exception as exc:
        logger.warning("OCR failed for %s: %s", path, exc)
        return ""

    return _finalize_text(text, max_chars)


def _ocr_page_images(
    images: Sequence[Any],
    pytesseract_module: Any,
    path: Path,
    page_num: int,
) -> list[str]:
    """Helper to run OCR on a list of images for a given page."""
    page_texts: list[str] = []
    for image in images:
        try:
            text = cast(str, pytesseract_module.image_to_string(image))
            if text and text.strip():
                page_texts.append(text.strip())
        except Exception as exc:
            logger.warning("OCR failed for %s page %d: %s", path, page_num, exc)
    return page_texts


def _extract_pdf_with_ocr(path: Path, max_chars: int | None) -> str:
    """Fallback OCR for PDFs that contain no embedded text."""
    try:
        from pdf2image import convert_from_path, pdfinfo_from_path  # type: ignore
    except ImportError:
        logger.info("pdf2image not installed; skipping PDF OCR for %s", path)
        return ""
    except Exception as exc:
        logger.warning("Failed to import pdf2image for %s: %s", path, exc)
        return ""

    try:
        import pytesseract  # type: ignore
    except ImportError:
        logger.info("pytesseract not installed; skipping PDF OCR")
        return ""
    except Exception as exc:
        logger.warning("Failed to import pytesseract for %s: %s", path, exc)
        return ""

    pytesseract_module = cast(Any, pytesseract)
    convert = cast(Callable[..., Sequence[Any]], convert_from_path)
    aggregated: list[str] = []

    try:
        info = pdfinfo_from_path(str(path))
        page_count = int(info.get("Pages", 0))
    except Exception as exc:
        logger.warning("Unable to read PDF page count for %s: %s", path, exc)
        page_count = 0

    if page_count > 0:
        for page_num in range(1, page_count + 1):
            try:
                images = convert(str(path), first_page=page_num, last_page=page_num)
                aggregated.extend(
                    _ocr_page_images(images, pytesseract_module, path, page_num)
                )
            except Exception as exc:
                logger.warning(
                    "Unable to rasterize PDF %s page %d for OCR: %s",
                    path,
                    page_num,
                    exc,
                )
    else:
        try:
            images = convert(str(path))
            aggregated.extend(_ocr_page_images(images, pytesseract_module, path, 0))
        except Exception as exc:
            logger.warning("Unable to rasterize PDF %s for OCR: %s", path, exc)

    if not aggregated:
        return ""

    return _finalize_text("\n".join(aggregated), max_chars)


def _extract_tables_with_camelot(path: Path) -> list[str]:
    """Extract tables from a PDF using the Camelot library."""
    tables_output: list[str] = []
    try:
        import camelot  # type: ignore

        camelot_module = cast(Any, camelot)
        camelot_tables = cast(
            Sequence[Any], camelot_module.read_pdf(str(path), pages="all")
        )
        for index, table in enumerate(camelot_tables):
            df = cast(Any, getattr(table, "df", None))
            csv_data = df.to_csv(index=False, header=True) if df is not None else ""
            tables_output.append(f"[Camelot Table {index + 1}]\n{csv_data}")
    except ImportError:
        logger.info("Camelot not installed; falling back to pdfplumber for %s", path)
    except Exception as exc:
        logger.warning("Camelot table extraction failed for %s: %s", path, exc)
    return tables_output


def _pdfplumber_table_to_csv(table: Sequence[Sequence[str | None]]) -> str:
    """Convert a single table extracted by pdfplumber to a CSV string."""
    csv_buffer = io.StringIO()
    writer = csv.writer(csv_buffer)
    for row in table:
        if not row or not any(row):
            continue
        normalized_row = [cell or "" for cell in row]
        writer.writerow(normalized_row)
    return csv_buffer.getvalue().strip()


def _extract_tables_with_pdfplumber(path: Path) -> list[str]:
    """Extract tables from a PDF using the pdfplumber library."""
    tables_output: list[str] = []
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        logger.info("pdfplumber not installed; skipping table extraction for %s", path)
        return []
    except Exception as exc:
        logger.warning("Failed to import pdfplumber for %s: %s", path, exc)
        return []

    try:
        pdfplumber_module = cast(Any, pdfplumber)
        with pdfplumber_module.open(str(path)) as pdf:
            pages = cast(Sequence[Any], getattr(pdf, "pages", []))
            for page_num, page in enumerate(pages, start=1):
                tables = cast(
                    Sequence[Sequence[Sequence[str | None]]] | None,
                    page.extract_tables(),
                )
                if not tables:
                    continue
                for table_index, table_data in enumerate(tables, start=1):
                    csv_data = _pdfplumber_table_to_csv(table_data)
                    if csv_data:
                        tables_output.append(
                            f"[pdfplumber Table {page_num}.{table_index}]\n{csv_data}"
                        )
    except Exception as exc:
        logger.warning("pdfplumber table extraction failed for %s: %s", path, exc)
    return tables_output


def _extract_pdf_tables(path: Path, max_chars: int | None) -> str:
    """Extract tables from PDFs using Camelot or pdfplumber when possible."""
    tables_output = _extract_tables_with_camelot(path)

    if not tables_output:
        tables_output = _extract_tables_with_pdfplumber(path)

    if not tables_output:
        return ""

    return _finalize_text("\n\n".join(tables_output), max_chars)


def _html_to_text_bs(html: str) -> str:
    """Convert HTML to text using BeautifulSoup."""
    from bs4 import BeautifulSoup  # type: ignore

    soup = cast(Any, BeautifulSoup(html, "html.parser"))
    # Remove script/style
    for tag in cast(Iterable[Any], soup(["script", "style", "noscript"])):
        tag.decompose()
    text = cast(str, soup.get_text(separator=" ", strip=True))
    return re.sub(r"\s+", " ", text)


def _html_to_text_regex(html: str) -> str:
    """Fallback HTML to text conversion using regex."""
    text = re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html)
    text = re.sub(r"<[^>]+>", " ", text)
    return re.sub(r"\s+", " ", text)


def _html_to_text(html: str) -> str:
    """Best-effort conversion of HTML to text; falls back to regex strip."""
    if not html:
        return ""
    try:
        return _html_to_text_bs(html)
    except ImportError:
        logger.info("BeautifulSoup not installed; falling back to regex for HTML.")
    except Exception as exc:
        logger.warning(
            "BeautifulSoup processing failed, falling back to regex: %s", exc
        )
    return _html_to_text_regex(html)


def _extract_text_from_doc_win32(path: Path) -> str:
    """Use pywin32/Word to extract text from legacy .doc files on Windows."""
    word = None
    doc = None
    try:
        import win32com.client

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(path.resolve()))
        return strip_control_chars(doc.Content.Text or "")
    except ImportError:
        logger.info("pywin32 not installed; cannot process .doc files on Windows.")
        return ""
    except Exception as e:
        logger.error("Error processing .doc file %s with win32com: %s", path, e)
        return ""
    finally:
        if doc:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word:
            try:
                word.Quit()
            except Exception:
                pass


def _extract_eml_headers(msg: Any) -> list[str]:
    """Extract key headers from an email message object."""
    headers: list[str] = []
    for hdr in ("From", "To", "Cc", "Bcc", "Subject", "Date"):
        if msg.get(hdr):
            headers.append(f"{hdr}: {msg.get(hdr)}")
    return headers


def _extract_eml_body(msg: Any) -> str:
    """Extract the text body from an email message object."""
    bodies: list[str] = []
    if msg.is_multipart():
        # First, try to find a text/plain part
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                with contextlib.suppress(Exception):
                    bodies.append(part.get_content())
        # If no plain text, fall back to HTML
        if not bodies:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    with contextlib.suppress(Exception):
                        bodies.append(_html_to_text(part.get_content()))
    else:
        # Not multipart, so get the single body
        try:
            if msg.get_content_type() == "text/html":
                bodies.append(_html_to_text(msg.get_content()))
            else:
                bodies.append(msg.get_content())
        except Exception:
            pass
    return "\n\n".join(bodies)


def _extract_eml(path: Path) -> str:
    """Parse .eml messages using the stdlib 'email' package."""
    try:
        from email import policy
        from email.parser import BytesParser

        with path.open("rb") as fh:
            msg = BytesParser(policy=policy.default).parse(fh)
    except Exception as e:
        logger.warning("Failed to parse EML %s: %s", path, e)
        return ""

    headers = _extract_eml_headers(msg)
    body = _extract_eml_body(msg)

    parts = headers + ["", body] if body else headers
    return strip_control_chars("\n".join(parts)).strip()


def _get_msg_body(msg_obj: Any) -> str:
    """Extracts the body from an extract_msg message object."""
    try:
        html_body = cast(str | None, getattr(msg_obj, "htmlBody", None))
        if html_body:
            return _html_to_text(html_body)
    except Exception:
        pass
    return cast(str | None, getattr(msg_obj, "body", "")) or ""


def _get_msg_headers(msg_obj: Any) -> list[str]:
    """Extracts headers from an extract_msg message object."""
    headers: list[str] = []
    for key in ("from", "to", "cc", "bcc", "subject", "date"):
        val = cast(str | None, getattr(msg_obj, key, None))
        if val:
            headers.append(f"{key.capitalize()}: {val}")
    return headers


def _close_msg_safely(msg_obj: Any, path: Path) -> None:
    """Safely close an extract_msg message object, trying multiple methods."""
    if msg_obj is None:
        return
    try:
        # First, try a standard .close() method
        close_candidate = cast(Any, msg_obj).close
        if callable(close_candidate):
            close_candidate()
            return
    except AttributeError:
        pass
    except Exception as close_error:
        logger.debug("Failed on .close() for MSG %s: %s", path, close_error)
        return  # Stop trying if .close() exists but fails

    try:
        # Fallback to context manager __exit__
        exit_candidate = cast(Any, msg_obj).__exit__
        if callable(exit_candidate):
            exit_candidate(None, None, None)
    except AttributeError:
        logger.debug("MSG object for %s has no close method.", path)
    except Exception as exit_error:
        logger.debug("Failed on .__exit__() for MSG %s: %s", path, exit_error)


def _extract_msg(path: Path) -> str:
    """Parse Outlook .msg files if extract_msg is available."""
    try:
        import extract_msg
    except ImportError:
        logger.info("extract_msg not installed; skipping .msg file: %s", path)
        return ""
    except Exception as exc:
        logger.warning("Failed to import extract_msg for %s: %s", path, exc)
        return ""

    msg_obj: Any | None = None
    try:
        msg_obj = extract_msg.Message(str(path))
        body = _get_msg_body(msg_obj)
        headers = _get_msg_headers(msg_obj)
        text = "\n".join([*headers, "", body])
        return strip_control_chars(text)
    except Exception as e:
        logger.warning("Failed to parse MSG %s: %s", path, e)
        return ""
    finally:
        _close_msg_safely(msg_obj, path)


def _extract_docx_content(path: Path, max_chars: int | None) -> str:
    """Extract content from a .docx file using python-docx."""
    import docx

    document = docx.Document(str(path))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        if paragraph.text and paragraph.text.strip():
            parts.append(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text and cell.text.strip():
                    parts.append(cell.text)

    text = "\n".join(parts)
    return _finalize_text(text, max_chars)


def _extract_doc_content(path: Path, max_chars: int | None) -> str:
    """Extract content from a legacy .doc file."""
    if os.name == "nt":
        win_text = _extract_text_from_doc_win32(path)
        if win_text:
            return _finalize_text(win_text, max_chars)

    # Cross-platform best-effort: try textract if installed
    try:
        import textract

        raw = textract.process(str(path))
        text = raw.decode("utf-8", errors="ignore")
        return _finalize_text(text, max_chars)
    except Exception:
        logger.info(
            "No supported reader for legacy .doc file on this platform: %s", path
        )
        return ""


def _extract_word_document(path: Path, suffix: str, max_chars: int | None) -> str:
    """Extract text from Word documents (.docx, .doc)."""
    try:
        if suffix == ".docx":
            return _extract_docx_content(path, max_chars)
        return _extract_doc_content(path, max_chars)
    except ImportError:
        logger.info("python-docx/textract not installed, skipping Word file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read Word document %s: %s", path, e)
        return ""


def _get_shape_text(shape: Any) -> str | None:
    """Extracts text from a PowerPoint shape if it exists."""
    if hasattr(shape, "text"):
        text = getattr(shape, "text", "") or ""
        return text.strip() if text.strip() else None
    return None


def _extract_powerpoint(path: Path, max_chars: int | None) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.info("python-pptx not installed, skipping PowerPoint file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to import python-pptx for %s: %s", path, e)
        return ""

    parts: list[str] = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                text = _get_shape_text(shape)
                if text:
                    parts.append(text)
    except Exception as e:
        logger.warning("Failed to read PowerPoint file %s: %s", path, e)
        return ""

    text = "\n".join(parts)
    return strip_control_chars(text[:max_chars] if max_chars else text)


def _extract_rtf(path: Path, max_chars: int | None) -> str:
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore

        data = path.read_bytes().decode("latin-1", errors="ignore")
        text = rtf_to_text(data)
        return strip_control_chars(text[:max_chars] if max_chars else text)
    except ImportError:
        logger.info("striprtf not installed, skipping RTF file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read RTF file %s: %s", path, e)
        return ""


def _extract_pdf(path: Path, max_chars: int | None) -> str:
    """Extract text from PDF files."""
    try:
        from pypdf import PdfReader  # type: ignore

        try:
            with path.open("rb") as fh:
                pdf = PdfReader(fh)
                # Try empty-password decryption when possible
                if getattr(pdf, "is_encrypted", False):
                    try:
                        pdf.decrypt("")  # type: ignore[attr-defined]
                    except Exception:
                        logger.warning(
                            "Skipping encrypted PDF (unable to decrypt): %s", path
                        )
                        return ""
                parts: list[str] = []
                acc = 0
                budget = max_chars if max_chars is not None else None
                for i, page in enumerate(getattr(pdf, "pages", [])):
                    try:
                        t = page.extract_text() or ""
                        if not t:
                            continue
                        remain = None if budget is None else (budget - acc)
                        if remain is not None and remain <= 0:
                            break
                        if remain is not None and len(t) > remain:
                            t = t[:remain]
                        parts.append(t)
                        acc += len(t)
                    except Exception as e:
                        logger.warning(
                            "Failed to extract text from page %d of %s: %s",
                            i + 1,
                            path,
                            e,
                        )
            text = "\n".join(parts)
            return strip_control_chars(text)
        except Exception as e:
            logger.warning("Failed to read PDF %s: %s. Skipping.", path, e)
            return ""
    except ImportError:
        logger.info("pypdf not installed, skipping PDF file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Unexpected error while reading PDF %s: %s", path, e)
        return ""


def _extract_excel(path: Path, suffix: str, max_chars: int | None) -> str:
    """Extract text from Excel files."""
    try:
        from typing import Literal

        import pandas as pd  # type: ignore

        engine: Literal["openpyxl", "xlrd"] = (
            "openpyxl" if suffix == ".xlsx" else "xlrd"
        )

        # Try explicit engine first, then fall back to pandas auto-detection
        def _iter_sheets(_engine: Literal["openpyxl", "xlrd"] | None):
            if _engine:
                return pd.ExcelFile(str(path), engine=_engine)
            return pd.ExcelFile(str(path))

        xl = None
        try:
            xl = _iter_sheets(engine)
        except Exception:
            xl = _iter_sheets(None)

        parts: list[str] = []
        acc = 0
        budget = max_chars if max_chars is not None else None
        max_cells = 200000  # Default max cells limit

        # Ensure the handle is closed even if exceptions occur
        with xl:
            for sheet_name in xl.sheet_names:
                try:
                    df = pd.read_excel(xl, sheet_name=sheet_name, dtype=str)
                    if df.size > max_cells and len(df.shape) > 1:
                        max_rows = max_cells // df.shape[1]
                        df = df.head(max_rows)
                    chunk = f"[Sheet: {sheet_name}]\n{df.to_csv(index=False)}"
                    remain = None if budget is None else (budget - acc)
                    if remain is not None and remain <= 0:
                        break
                    if remain is not None and len(chunk) > remain:
                        chunk = chunk[:remain]
                    parts.append(chunk)
                    acc += len(chunk)
                except Exception as e:
                    logger.warning(
                        "Failed reading sheet '%s' in %s: %s", sheet_name, path, e
                    )
        text = "\n".join(parts)
        return strip_control_chars(text)
    except ImportError:
        logger.info("pandas/openpyxl/xlrd not installed, skipping Excel file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read Excel file %s: %s", path, e)
        return ""


def extract_text(path: Path, *, max_chars: int | None = None) -> str:
    """
    Extract text from supported file types with robust error handling.

    Supports: .txt, .pdf, .docx, .doc, .xlsx, .xls, .pptx, .ppt,
    .rtf, .eml, .msg, .html, .xml, .md, .json, .yaml, .csv
    Unknown/binary formats return empty string.

    Args:
        path: Path to the file (must exist and be readable)
        max_chars: Optional hard cap on returned text size

    Returns:
        Extracted and sanitized text, possibly truncated.
        Empty string on errors or unsupported formats.
    """
    global _extraction_cache

    # 1. Cache Lookup
    cached = _check_cache(path, max_chars)
    if cached is not None:
        return cached

    # 2. Validation
    if not _validate_path(path):
        return ""
    path = path.resolve()
    suffix = path.suffix.lower()

    # 3. Text Extraction
    segments: list[str] = []
    seen_hashes: set[int] = set()

    if suffix in TEXT_EXTENSIONS:
        _extract_text_file(path, suffix, max_chars, segments, seen_hashes)
    else:
        # Build and run pipeline
        pipeline = _build_pipeline(path, suffix, max_chars)
        _run_pipeline(pipeline, segments, seen_hashes)

        # Fallbacks (PDF OCR, Tables, Images)
        _apply_fallbacks(path, suffix, max_chars, segments, seen_hashes)

    # 4. Finalize & Cache
    combined = "\n\n".join(segments)
    result = _finalize_text(combined, max_chars) if combined else ""
    _update_cache(path, max_chars, result)

    return result


def _check_cache(path: Path, max_chars: int | None) -> str | None:
    try:
        cache_key = (path.resolve(), max_chars)
        with _extraction_cache_lock:
            if cache_key in _extraction_cache:
                entry = _extraction_cache[cache_key]
                if _is_cache_valid(entry, path=path):
                    logger.debug("Using cached text extraction for %s", path.name)
                    return entry[3]
    except Exception as e:
        logger.debug("Cache lookup failed for %s: %s", path, e)
    return None


def _validate_path(path: Path) -> bool:
    try:
        if not path.exists():
            logger.debug("Path does not exist: %s", path)
            return False
        if not path.is_file():
            logger.debug("Path is not a file: %s", path)
            return False
    except (ValueError, OSError) as e:
        logger.warning("Invalid path: %s - %s", path, e)
        return False
    return True


def _add_segment(segment: str, segments: list[str], seen_hashes: set[int]) -> None:
    if not segment:
        return
    # Strip control characters before hashing to ensure that segments differing
    # only by invisible characters are treated as duplicates.
    normalized = strip_control_chars(segment).strip()
    if not normalized:
        return
    signature = hash(normalized)
    if signature in seen_hashes:
        return
    seen_hashes.add(signature)
    segments.append(normalized)


def _extract_text_file(
    path: Path,
    suffix: str,
    max_chars: int | None,
    segments: list[str],
    seen_hashes: set[int],
) -> None:
    from cortex.utils import read_text_file

    try:
        text = read_text_file(path, max_chars=max_chars)
        if suffix in {".html", ".htm", ".xml"}:
            text = _html_to_text(text)
        _add_segment(text, segments, seen_hashes)
    except OSError as e:
        logger.warning("Failed to read text file %s: %s", path, e)


def _build_pipeline(
    path: Path, suffix: str, max_chars: int | None
) -> list[tuple[Callable[[], str], bool]]:
    pipeline: list[tuple[Callable[[], str], bool]] = []

    # Helper lambdas
    def uns():
        return _extract_with_unstructured(path, suffix, max_chars)

    def tika():
        return _extract_with_tika(path, max_chars)

    if suffix in EMAIL_EXTENSIONS:
        if suffix == ".eml":
            pipeline.append((lambda: _extract_eml(path), False))
        else:
            pipeline.append((lambda: _extract_msg(path), False))
        pipeline.append((uns, True))
        pipeline.append((tika, True))

    elif suffix in DOCX_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((lambda: _extract_word_document(path, suffix, max_chars), True))
        pipeline.append((tika, False))

    elif suffix in PDF_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((tika, True))
        pipeline.append((lambda: _extract_pdf(path, max_chars), True))

    elif suffix in EXCEL_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((lambda: _extract_excel(path, suffix, max_chars), True))
        pipeline.append((tika, False))

    elif suffix in PPT_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((lambda: _extract_powerpoint(path, max_chars), True))
        pipeline.append((tika, False))

    elif suffix in RTF_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((lambda: _extract_rtf(path, max_chars), True))
        pipeline.append((tika, False))

    elif suffix in IMAGE_EXTENSIONS:
        pipeline.append((uns, False))
        pipeline.append((tika, False))
        pipeline.append((lambda: _extract_image_ocr(path, max_chars), True))

    elif suffix in MBOX_EXTENSIONS:
        pipeline.append((tika, False))
        pipeline.append((uns, True))

    else:
        pipeline.append((uns, False))
        pipeline.append((tika, False))

    return pipeline


def _run_pipeline(
    pipeline: list[tuple[Callable[[], str], bool]],
    segments: list[str],
    seen_hashes: set[int],
) -> None:
    for extractor, always_run in pipeline:
        if segments and not always_run:
            continue
        segment = extractor()
        if segment:
            _add_segment(segment, segments, seen_hashes)


def _apply_fallbacks(
    path: Path,
    suffix: str,
    max_chars: int | None,
    segments: list[str],
    seen_hashes: set[int],
) -> None:
    if suffix in PDF_EXTENSIONS:
        if not segments:
            ocr_text = _extract_pdf_with_ocr(path, max_chars)
            if ocr_text:
                _add_segment(ocr_text, segments, seen_hashes)
        table_text = _extract_pdf_tables(path, max_chars)
        if table_text:
            _add_segment(table_text, segments, seen_hashes)
    elif suffix in IMAGE_EXTENSIONS and not segments:
        ocr_text = _extract_image_ocr(path, max_chars)
        if ocr_text:
            _add_segment(ocr_text, segments, seen_hashes)


def _update_cache(path: Path, max_chars: int | None, result: str) -> None:
    global _extraction_cache
    if not result:
        return

    try:
        cache_key = (path.resolve(), max_chars)
        with _extraction_cache_lock:
            try:
                st = path.stat()
                _extraction_cache[cache_key] = (
                    time.time(),
                    float(st.st_mtime),
                    int(st.st_size),
                    result,
                )
            except Exception:
                _extraction_cache[cache_key] = (time.time(), 0.0, 0, result)

            # Clean old cache entries periodically
            if len(_extraction_cache) > 100:
                _extraction_cache = {
                    k: v
                    for k, v in _extraction_cache.items()
                    if (time.time() - float(v[0])) <= _CACHE_TTL
                }
    except Exception as e:
        logger.debug("Failed to update cache for %s: %s", path, e)
