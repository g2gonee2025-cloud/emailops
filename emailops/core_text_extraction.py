from __future__ import annotations

import asyncio
import contextlib
import logging
import os
import re
import threading
import time
from functools import partial
from pathlib import Path

from .util_files import _strip_control_chars, read_text_file

try:
    from .core_config import EmailOpsConfig
except ImportError:
    from core_config import EmailOpsConfig  # type: ignore

"""
Text extraction utilities for various file formats.
Handles PDF, Word, Excel, PowerPoint, RTF, EML, and MSG files.

Note: For any JSON extraction, use `scrub_json_string` from utils before parsing to ensure control characters are removed.
"""

logger = logging.getLogger(__name__)

# File type support
TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".log",
    ".json",
    ".yaml",
    ".yml",
    ".csv",
    ".xml",
    ".html",
    ".htm",
}
DOCX_EXTENSIONS = {".docx", ".doc"}
PDF_EXTENSIONS = {".pdf"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
PPT_EXTENSIONS = {".pptx", ".ppt"}
RTF_EXTENSIONS = {".rtf"}
EMAIL_EXTENSIONS = {".eml", ".msg"}

# Global cache for expensive text extraction operations
_extraction_cache: dict[tuple[Path, int | None], tuple[float, str]] = {}
_extraction_cache_lock = threading.Lock()
_CACHE_TTL = 3600  # 1 hour TTL for cache entries


def _is_cache_valid(cache_entry: tuple[float, str]) -> bool:
    """Check if cache entry is still valid based on TTL."""
    cache_time, _ = cache_entry
    return (time.time() - cache_time) < _CACHE_TTL


def _html_to_text(html: str) -> str:
    """Best-effort conversion of HTML to text; falls back to regex strip."""
    if not html:
        return ""
    # Try BeautifulSoup if available for better results
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(html, "html.parser")
        # Remove script/style
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator=" ", strip=True)
        return re.sub(r"\s+", " ", text)
    except Exception:
        # Regex fallback: strip tags and collapse whitespace
        text = re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html)
        text = re.sub(r"<[^>]+>", " ", text)
        return re.sub(r"\s+", " ", text)


def _extract_text_from_doc_win32(path: Path) -> str:
    """Use pywin32/Word to extract text from legacy .doc files on Windows."""
    try:
        import win32com.client

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = None
        try:
            doc = word.Documents.Open(str(path.resolve()))
            return _strip_control_chars(doc.Content.Text or "")
        finally:
            if doc:
                doc.Close(False)
            word.Quit()
    except ImportError:
        logger.info("pywin32 not installed; cannot process .doc files on Windows.")
        return ""
    except Exception as e:
        logger.error("Error processing .doc file %s with win32com: %s", path, e)
        return ""


def _extract_eml(path: Path) -> str:
    """Parse .eml messages using the stdlib 'email' package."""
    try:
        from email import policy
        from email.parser import BytesParser

        msg = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
    except Exception as e:
        logger.warning("Failed to parse EML %s: %s", path, e)
        return ""

    parts: list[str] = []
    # Include a minimal header block for context
    for hdr in ("From", "To", "Cc", "Bcc", "Subject", "Date"):
        if msg.get(hdr):
            parts.append(f"{hdr}: {msg.get(hdr)}")

    # Prefer text/plain; fall back to text/html
    bodies: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            if ctype == "text/plain":
                with contextlib.suppress(Exception):
                    bodies.append(part.get_content())
        if not bodies:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    with contextlib.suppress(Exception):
                        bodies.append(_html_to_text(part.get_content()))
    else:
        ctype = msg.get_content_type()
        try:
            if ctype == "text/html":
                bodies.append(_html_to_text(msg.get_content()))
            else:
                bodies.append(msg.get_content())
        except Exception:
            pass

    parts.append("")  # blank line between headers and body
    parts.append("\n\n".join(bodies))
    return _strip_control_chars("\n".join(parts)).strip()


def _extract_msg(path: Path) -> str:
    """Parse Outlook .msg files if extract_msg is available."""
    try:
        import extract_msg  # type: ignore
    except Exception:
        logger.info("extract_msg not installed; skipping .msg file: %s", path)
        return ""

    m = None
    text_result = ""

    try:
        m = extract_msg.Message(str(path))
        # Prefer HTML if available
        body = ""
        try:
            html = getattr(m, "htmlBody", None)
            if html:
                body = _html_to_text(html)
        except Exception:
            body = ""

        if not body:
            body = m.body or ""
        headers = []
        for k in ("from", "to", "cc", "bcc", "subject", "date"):
            val = getattr(m, k, None)
            if val:
                headers.append(f"{k.capitalize()}: {val}")

        text = "\n".join([*headers, "", body])
        text_result = _strip_control_chars(text)
    except Exception as e:
        logger.warning("Failed to parse MSG %s: %s", path, e)
        text_result = ""
    finally:
        # HIGH #8: Improved resource cleanup with explicit error handling
        # Close the message handle to prevent resource leaks
        if m is not None:
            try:
                if hasattr(m, 'close') and callable(getattr(m, 'close', None)):
                    m.close()  # type: ignore[attr-defined]
                elif hasattr(m, '__exit__'):
                    # If it's a context manager, try to exit it
                    m.__exit__(None, None, None)  # type: ignore[attr-defined]
            except Exception as close_error:
                # Log close errors but don't fail the extraction
                logger.debug("Failed to close MSG handle for %s: %s", path, close_error)

    return text_result


def extract_text(path: Path, *, max_chars: int | None = None, use_cache: bool = True) -> str:
    """
    Extract text from supported file types with robust error handling and caching.

    Supports: .txt, .pdf, .docx, .doc, .xlsx, .xls, .pptx, .ppt,
    .rtf, .eml, .msg, .html, .xml, .md, .json, .yaml, .csv
    Unknown/binary formats return empty string.

    Args:
        path: Path to the file (must exist and be readable)
        max_chars: Optional hard cap on returned text size
        use_cache: Whether to use cached results if available

    Returns:
        Extracted and sanitized text, possibly truncated.
        Empty string on errors or unsupported formats.
    """
    # Declare global variable for modification
    global _extraction_cache

    # Check cache first
    if use_cache:
        cache_key = (path.resolve(), max_chars)
        with _extraction_cache_lock:
            if cache_key in _extraction_cache:
                cached_data = _extraction_cache[cache_key]
                if _is_cache_valid(cached_data):
                    logger.debug("Using cached text extraction for %s", path)
                    return cached_data[1]

    # Basic path validation
    try:
        if not path.exists():
            logger.debug("Path does not exist: %s", path)
            return ""
        if not path.is_file():
            logger.debug("Path is not a file: %s", path)
            return ""
        path = path.resolve()
    except (ValueError, OSError) as e:
        logger.warning("Invalid path: %s - %s", path, e)
        return ""

    cfg = EmailOpsConfig.load()
    if not any(path.match(pattern) for pattern in cfg.allowed_file_patterns):
        logger.debug("Skipping file with disallowed extension: %s", path)
        return ""

    suffix = path.suffix.lower()
    result = ""

    # Text-like files (handle HTML/XML below)
    if suffix in TEXT_EXTENSIONS:
        text = read_text_file(path, max_chars=max_chars)
        if suffix in {".html", ".htm", ".xml"}:
            text = _html_to_text(text)
        result = text

    # Word documents
    elif suffix in DOCX_EXTENSIONS:
        result = _extract_word_document(path, suffix, max_chars)

    # PowerPoint
    elif suffix in PPT_EXTENSIONS:
        result = _extract_powerpoint(path, max_chars)

    # RTF
    elif suffix in RTF_EXTENSIONS:
        result = _extract_rtf(path, max_chars)

    # PDFs
    elif suffix in PDF_EXTENSIONS:
        result = _extract_pdf(path, max_chars)

    # Excel
    elif suffix in EXCEL_EXTENSIONS:
        result = _extract_excel(path, suffix, max_chars)

    # Raw email files
    elif suffix in EMAIL_EXTENSIONS:
        result = _extract_eml(path) if suffix == ".eml" else _extract_msg(path)

    else:
        logger.debug("Unsupported file format for text extraction: %s", suffix)

    # Cache the result if caching is enabled
    if use_cache and result:
        cache_key = (path.resolve(), max_chars)
        with _extraction_cache_lock:
            _extraction_cache[cache_key] = (time.time(), result)

            # Clean old cache entries periodically
            if len(_extraction_cache) > 100:  # Clean when cache gets too large
                _extraction_cache = {k: v for k, v in _extraction_cache.items() if _is_cache_valid(v)}

    return result


def _extract_word_document(path: Path, suffix: str, max_chars: int | None) -> str:
    """Extract text from Word documents (.docx, .doc)."""
    try:
        if suffix == ".docx":
            import docx  # type: ignore

            doc = docx.Document(str(path))
            parts: list[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            # Include table cells
            for table in getattr(doc, "tables", []):
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            parts.append(cell.text)
            text = "\n".join(parts)
            if max_chars is not None and len(text) > max_chars:
                text = text[:max_chars]
            return _strip_control_chars(text)
        else:  # .doc
            if os.name == "nt":
                win_text = _extract_text_from_doc_win32(path)
                if win_text:
                    return win_text[:max_chars] if max_chars else win_text
            # Cross-platform best-effort: try textract if installed
            try:
                import textract  # type: ignore

                raw = textract.process(str(path))  # bytes
                txt = raw.decode("utf-8", errors="ignore")
                return _strip_control_chars(txt[:max_chars] if max_chars else txt)
            except Exception:
                logger.info("No supported reader for legacy .doc file on this platform: %s", path)
                return ""
    except ImportError:
        logger.info("python-docx/textract not installed, skipping Word file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read Word document %s: %s", path, e)
        return ""


def _extract_powerpoint(path: Path, max_chars: int | None) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # Many shapes expose .text; ignore drawing objects without text
                if hasattr(shape, "text"):
                    t = getattr(shape, "text", "") or ""
                    if t.strip():
                        parts.append(t)
        text = "\n".join(parts)
        return _strip_control_chars(text[:max_chars] if max_chars else text)
    except ImportError:
        logger.info("python-pptx not installed, skipping PowerPoint file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read PowerPoint file %s: %s", path, e)
        return ""


def _extract_rtf(path: Path, max_chars: int | None) -> str:
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore

        data = path.read_bytes().decode("latin-1", errors="ignore")
        text = rtf_to_text(data)
        return _strip_control_chars(text[:max_chars] if max_chars else text)
    except ImportError:
        logger.info("striprtf not installed, skipping RTF file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read RTF file %s: %s", path, e)
        return ""


def _extract_pdf(path: Path, max_chars: int | None) -> str:
    """Extract text from PDF files."""
    try:
        from pypdf import PdfReader  # type: ignore

        try:
            with Path.open(path, "rb") as fh:
                pdf = PdfReader(fh)
                # Try empty-password decryption when possible
                if getattr(pdf, "is_encrypted", False):
                    try:
                        pdf.decrypt("")  # type: ignore[attr-defined]
                    except Exception:
                        logger.warning("Skipping encrypted PDF (unable to decrypt): %s", path)
                        return ""
                parts: list[str] = []
                acc = 0
                budget = max_chars if max_chars is not None else None
                for i, page in enumerate(getattr(pdf, "pages", [])):
                    try:
                        t = page.extract_text() or ""
                        if not t:
                            continue
                        remain = None if budget is None else (budget - acc)
                        if remain is not None and remain <= 0:
                            break
                        if remain is not None and len(t) > remain:
                            t = t[:remain]
                        parts.append(t)
                        acc += len(t)
                    except Exception as e:
                        logger.warning("Failed to extract text from page %d of %s: %s", i + 1, path, e)
            text = "\n".join(parts)
            return _strip_control_chars(text)
        except Exception as e:
            logger.warning("Failed to read PDF %s: %s. Skipping.", path, e)
            return ""
    except ImportError:
        logger.info("pypdf not installed, skipping PDF file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Unexpected error while reading PDF %s: %s", path, e)
        return ""


def _extract_excel(path: Path, suffix: str, max_chars: int | None) -> str:
    """Extract text from Excel files."""
    try:
        from typing import Literal

        import pandas as pd  # type: ignore

        engine: Literal["openpyxl", "xlrd"] = "openpyxl" if suffix == ".xlsx" else "xlrd"

        # Try explicit engine first, then fall back to pandas auto-detection
        def _iter_sheets(_engine: Literal["openpyxl", "xlrd"] | None):
            if _engine:
                return pd.ExcelFile(str(path), engine=_engine)
            return pd.ExcelFile(str(path))

        xl = None
        try:
            xl = _iter_sheets(engine)
        except Exception:
            xl = _iter_sheets(None)

        parts: list[str] = []
        acc = 0
        budget = max_chars if max_chars is not None else None
        max_cells = 200000  # Default max cells limit

        # Ensure the handle is closed even if exceptions occur
        with xl:
            for sheet_name in xl.sheet_names:
                try:
                    df = pd.read_excel(xl, sheet_name=sheet_name, dtype=str)
                    if df.size > max_cells and len(df.shape) > 1:
                        max_rows = max_cells // df.shape[1]
                        df = df.head(max_rows)
                    chunk = f"[Sheet: {sheet_name}]\n{df.to_csv(index=False)}"
                    remain = None if budget is None else (budget - acc)
                    if remain is not None and remain <= 0:
                        break
                    if remain is not None and len(chunk) > remain:
                        chunk = chunk[:remain]
                    parts.append(chunk)
                    acc += len(chunk)
                except Exception as e:
                    logger.warning("Failed reading sheet '%s' in %s: %s", sheet_name, path, e)
        text = "\n".join(parts)
        return _strip_control_chars(text)
    except ImportError:
        logger.info("pandas/openpyxl/xlrd not installed, skipping Excel file: %s", path)
        return ""
    except Exception as e:
        logger.warning("Failed to read Excel file %s: %s", path, e)
        return ""


async def extract_text_async(path: Path, *, max_chars: int | None = None) -> str:
    """
    Async version of extract_text using thread pool.

    Args:
        path: Path to the file
        max_chars: Optional hard cap on returned text size

    Returns:
        Extracted and sanitized text
    """
    loop = asyncio.get_event_loop()
    func = partial(extract_text, path, max_chars=max_chars)
    return await loop.run_in_executor(None, func)
